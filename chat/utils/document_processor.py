import PyPDF2
import pandas as pd
from docx import Document as DocxDocument
from pptx import Presentation
import os

class DocumentProcessor:
    @staticmethod
    def process_document(file_path):
        """Process different types of documents and extract text content."""
        extension = os.path.splitext(file_path)[1].lower()
        
        if extension == '.pdf':
            return DocumentProcessor._process_pdf(file_path)
        elif extension == '.docx':
            return DocumentProcessor._process_docx(file_path)
        elif extension == '.xlsx' or extension == '.xls':
            return DocumentProcessor._process_excel(file_path)
        elif extension == '.csv':
            return DocumentProcessor._process_csv(file_path)
        elif extension == '.pptx':
            return DocumentProcessor._process_pptx(file_path)
        else:
            raise ValueError(f"Unsupported file type: {extension}")

    @staticmethod
    def _process_pdf(file_path):
        text = []
        with open(file_path, 'rb') as file:
            pdf_reader = PyPDF2.PdfReader(file)
            for page in pdf_reader.pages:
                text.append(page.extract_text())
        return '\n'.join(text)

    @staticmethod
    def _process_docx(file_path):
        doc = DocxDocument(file_path)
        return '\n'.join([paragraph.text for paragraph in doc.paragraphs])

    @staticmethod
    def _process_excel(file_path):
        df = pd.read_excel(file_path)
        return df.to_string()

    @staticmethod
    def _process_csv(file_path):
        df = pd.read_csv(file_path)
        return df.to_string()

    @staticmethod
    def _process_pptx(file_path):
        prs = Presentation(file_path)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return '\n'.join(text)

def create_chunks(text, chunk_size=1000):
    """Split text into chunks of approximately equal size."""
    words = text.split()
    chunks = []
    current_chunk = []
    current_size = 0
    
    for word in words:
        current_size += len(word) + 1  # +1 for space
        if current_size > chunk_size:
            chunks.append(' '.join(current_chunk))
            current_chunk = [word]
            current_size = len(word)
        else:
            current_chunk.append(word)
    
    if current_chunk:
        chunks.append(' '.join(current_chunk))
    
    return chunks